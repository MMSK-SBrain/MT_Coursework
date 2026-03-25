#!/usr/bin/env python3
"""
Custom Markdown to PowerPoint Converter using python-pptx
Designed for automotive lecture frameworks with full styling control
"""

import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Color Scheme - Engineering Blue/Gray
COLORS = {
    'primary_blue': RGBColor(68, 114, 196),     # #4472C4
    'dark_blue': RGBColor(46, 80, 144),         # #2E5090
    'text_gray': RGBColor(51, 51, 51),          # #333333
    'light_gray': RGBColor(90, 90, 90),         # #5A5A5A
    'accent_orange': RGBColor(255, 102, 0),     # #FF6600
    'white': RGBColor(255, 255, 255),           # #FFFFFF
}

def clean_markdown_text(text):
    """Remove HTML comments and clean markdown syntax"""
    # Remove HTML comments
    text = re.sub(r'<!--.*?-->', '', text, flags=re.DOTALL)

    # Remove markdown bold/italic but keep the text
    text = re.sub(r'\*\*\*(.+?)\*\*\*', r'\1', text)  # bold+italic
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)      # bold
    text = re.sub(r'\*(.+?)\*', r'\1', text)          # italic
    text = re.sub(r'__(.+?)__', r'\1', text)          # bold (alt)
    text = re.sub(r'_(.+?)_', r'\1', text)            # italic (alt)

    # Clean up extra whitespace
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
    text = text.strip()

    return text

def parse_markdown_slides(markdown_content):
    """Parse markdown into slide data structures"""
    # Split by slide separators (---)
    raw_slides = re.split(r'\n---+\n', markdown_content)

    slides = []

    for raw_slide in raw_slides:
        raw_slide = raw_slide.strip()
        if not raw_slide or len(raw_slide) < 10:
            continue

        # Skip frontmatter (YAML between --- at start)
        if raw_slide.startswith('marp:') or raw_slide.startswith('theme:'):
            continue

        slide_data = parse_slide_content(raw_slide)
        if slide_data:
            slides.append(slide_data)

    return slides

def parse_slide_content(content):
    """Parse individual slide content into structured format"""
    lines = content.split('\n')

    slide = {
        'title': '',
        'subtitle': '',
        'content': [],
        'is_title_slide': False
    }

    i = 0
    while i < len(lines):
        line = lines[i].strip()

        # Skip empty lines and HTML comments
        if not line or line.startswith('<!--'):
            i += 1
            continue

        # H1 heading (# Title)
        if line.startswith('# '):
            slide['title'] = clean_markdown_text(line[2:])
            slide['is_title_slide'] = True
            i += 1
            continue

        # H2 heading (## Title)
        if line.startswith('## '):
            slide['title'] = clean_markdown_text(line[3:])
            i += 1
            continue

        # H3 heading (### Subtitle)
        if line.startswith('### '):
            slide['subtitle'] = clean_markdown_text(line[4:])
            i += 1
            continue

        # Bullet point (- item or * item)
        if line.startswith('- ') or line.startswith('* '):
            bullet_text = clean_markdown_text(line[2:])
            slide['content'].append({
                'type': 'bullet',
                'text': bullet_text,
                'level': 0
            })
            i += 1
            continue

        # Numbered list (1. item)
        if re.match(r'^\d+\.\s+', line):
            bullet_text = clean_markdown_text(re.sub(r'^\d+\.\s+', '', line))
            slide['content'].append({
                'type': 'bullet',
                'text': bullet_text,
                'level': 0
            })
            i += 1
            continue

        # Code block
        if line.startswith('```'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            slide['content'].append({
                'type': 'code',
                'text': '\n'.join(code_lines)
            })
            i += 1
            continue

        # Regular paragraph
        if line:
            slide['content'].append({
                'type': 'text',
                'text': clean_markdown_text(line)
            })

        i += 1

    return slide if slide['title'] or slide['content'] else None

def create_title_slide(prs, slide_data):
    """Create a title slide with blue gradient background"""
    slide_layout = prs.slide_layouts[0]  # Title layout
    slide = prs.slides.add_slide(slide_layout)

    # Set background to blue
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary_blue']

    # Title
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = slide_data['title']
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = COLORS['white']

    # Subtitle from first content item if exists
    if len(slide.placeholders) > 1 and slide_data['content']:
        subtitle = slide.placeholders[1]
        subtitle_text = '\n'.join([item['text'] for item in slide_data['content'] if item['type'] == 'text'][:3])
        subtitle.text = subtitle_text
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = COLORS['white']

    return slide

def create_content_slide(prs, slide_data):
    """Create a content slide with title and bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = slide_data['title']
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary_blue']

    # Content
    if len(slide.placeholders) > 1:
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        # Add subtitle if exists
        if slide_data['subtitle']:
            p = text_frame.paragraphs[0]
            p.text = slide_data['subtitle']
            p.level = 0
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = COLORS['dark_blue']

        # Add content items
        for item in slide_data['content']:
            if item['type'] == 'bullet':
                p = text_frame.add_paragraph() if text_frame.paragraphs[0].text else text_frame.paragraphs[0]
                p.text = item['text']
                p.level = item.get('level', 0)
                p.font.size = Pt(18)
                p.font.color.rgb = COLORS['text_gray']
                p.space_after = Pt(6)

            elif item['type'] == 'text':
                p = text_frame.add_paragraph() if text_frame.paragraphs[0].text else text_frame.paragraphs[0]
                p.text = item['text']
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = COLORS['text_gray']
                p.space_after = Pt(10)

            elif item['type'] == 'code':
                p = text_frame.add_paragraph() if text_frame.paragraphs[0].text else text_frame.paragraphs[0]
                p.text = item['text']
                p.level = 0
                p.font.name = 'Consolas'
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(120, 120, 120)

    return slide

def create_presentation(slides_data, output_path):
    """Create PowerPoint presentation from parsed slide data"""
    prs = Presentation()

    # Set slide dimensions (16:9 - standard widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_data:
        if slide_data['is_title_slide']:
            create_title_slide(prs, slide_data)
        else:
            create_content_slide(prs, slide_data)

    prs.save(output_path)
    return len(slides_data)

def convert_markdown_to_pptx(md_file, output_file):
    """Main conversion function"""
    print(f"Converting: {md_file.name}")

    try:
        # Read markdown
        with open(md_file, 'r', encoding='utf-8') as f:
            markdown_content = f.read()

        # Parse slides
        slides_data = parse_markdown_slides(markdown_content)

        if not slides_data:
            print(f"✗ No slides found in {md_file.name}")
            return False

        # Create PowerPoint
        slide_count = create_presentation(slides_data, output_file)

        print(f"✓ Created: {output_file}")
        print(f"  Slides: {slide_count}")

        return True

    except Exception as e:
        print(f"✗ Failed: {md_file.name}")
        print(f"  Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python markdown-to-pptx.py <input-markdown.md> [output-dir]")
        print("  If output-dir is not specified, uses ./pptx/")
        sys.exit(1)

    input_path = Path(sys.argv[1])

    if not input_path.exists():
        print(f"Error: File not found: {input_path}")
        sys.exit(1)

    # Output directory
    if len(sys.argv) >= 3:
        # Custom output directory provided
        output_dir = Path(sys.argv[2])
    else:
        # Default: current directory / pptx
        output_dir = Path.cwd() / 'pptx'

    output_dir.mkdir(parents=True, exist_ok=True)

    # Output file
    output_path = output_dir / (input_path.stem.replace('-Framework', '') + '.pptx')

    # Convert
    success = convert_markdown_to_pptx(input_path, output_path)

    sys.exit(0 if success else 1)
